"""
propale_core.py — Logique métier extraite de propale_depuis_excel_v3.py
Utilisé par app.py (interface Streamlit)
"""

import os
import re
import glob
import json
import shutil
import zipfile
import tempfile
import math as _math
from datetime import datetime
from openai import OpenAI

# ============================================================
# === CONFIGURATION ==========================================
# ============================================================

GROQ_MODEL = "llama-3.3-70b-versatile"

CUSTOM_PROMPTS_PATH = os.path.join(os.path.dirname(__file__), "custom_prompts.json")


def charger_prompts_personnalises() -> dict:
    if os.path.exists(CUSTOM_PROMPTS_PATH):
        try:
            with open(CUSTOM_PROMPTS_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return {}
    return {}


def sauvegarder_prompts_personnalises(custom: dict):
    with open(CUSTOM_PROMPTS_PATH, "w", encoding="utf-8") as f:
        json.dump(custom, f, ensure_ascii=False, indent=2)


def get_balises_effectives() -> dict:
    """
    Retourne BALISES_IA fusionné avec les prompts personnalisés.
    - Si la clé existe dans BALISES_IA  → surcharge le prompt par défaut.
    - Si la clé n'existe PAS dans BALISES_IA → nouvelle balise ajoutée.
    """
    custom = charger_prompts_personnalises()
    result = dict(BALISES_IA)
    result.update(custom)   # override + nouvelles balises
    return result


def est_balise_custom(key: str) -> bool:
    """Retourne True si la balise est une création personnalisée (absente de BALISES_IA)."""
    return key not in BALISES_IA

SLIDES_PAR_OPTION = {
    2: ["slide45.xml", "slide46.xml", "slide47.xml"],
    3: ["slide48.xml", "slide49.xml", "slide50.xml"],
}

# Groupes de slides optionnels (indices 0-based dans le template)
# Modifier les listes "slides" selon la numérotation réelle de votre template.
GROUPES_SLIDES = {
    "qualitatif": {
        "label": "Étude qualitative (ESD)",
        "description": "Entretiens semi-directifs — slides 34 à 37",
        "slides": [33, 34, 35, 36],
        "default": True,
    },
    "etude_mystere": {
        "label": "Étude mystère",
        "description": "Visites mystère — slides 38 et 39",
        "slides": [37, 38],
        "default": True,
    },
    "ideation": {
        "label": "Séance d'idéation",
        "description": "Atelier idéation — slides 40 et 41",
        "slides": [39, 40],
        "default": True,
    },
}

BOLD_TAGS  = {"{{sous_objun}}", "{{sous_objdeux}}", "{{sous_objtrois}}"}
UPPER_TAGS = {"{{objectif_principal}}"}

TAGS_FINANCIERS_PAR_OPTION = []


# ============================================================
# === CONVERSION MONTANT EN TOUTES LETTRES (FRANÇAIS) ========
# ============================================================

def montant_en_lettres(valeur: float) -> str:
    UNITES = ['', 'un', 'deux', 'trois', 'quatre', 'cinq', 'six', 'sept', 'huit', 'neuf',
              'dix', 'onze', 'douze', 'treize', 'quatorze', 'quinze', 'seize',
              'dix-sept', 'dix-huit', 'dix-neuf']
    DIZAINES = ['', 'dix', 'vingt', 'trente', 'quarante', 'cinquante',
                'soixante', 'soixante', 'quatre-vingt', 'quatre-vingt']

    def deux_chiffres(n):
        if n < 20:
            return UNITES[n]
        d, u = divmod(n, 10)
        if d == 7:
            u += 10
            lien = '-et-' if u == 11 else '-'
            return 'soixante' + (lien + UNITES[u] if u > 0 else '')
        if d == 9:
            u += 10
            return 'quatre-vingt' + ('-' + UNITES[u] if u > 0 else 's')
        if d == 8:
            return 'quatre-vingts' if u == 0 else 'quatre-vingt-' + UNITES[u]
        if u == 1:
            return DIZAINES[d] + '-et-un'
        if u == 0:
            return DIZAINES[d]
        return DIZAINES[d] + '-' + UNITES[u]

    def trois_chiffres(n):
        c, reste = divmod(n, 100)
        if c == 0:
            return deux_chiffres(reste)
        base = ('cent' if c == 1 else UNITES[c] + ' cent')
        if reste == 0:
            return base + ('' if c == 1 else 's')
        return base + ' ' + deux_chiffres(reste)

    euros    = int(valeur)
    centimes = round((valeur - euros) * 100)
    if centimes == 100:
        euros   += 1
        centimes = 0

    n = euros
    millions, reste = divmod(n, 1_000_000)
    milliers, unite = divmod(reste, 1000)

    parts = []
    if millions > 0:
        parts.append('un million' if millions == 1 else trois_chiffres(millions) + ' millions')
    if milliers > 0:
        parts.append('mille' if milliers == 1 else trois_chiffres(milliers) + ' mille')
    if unite > 0:
        parts.append(trois_chiffres(unite))

    result = 'zéro euro' if n == 0 else ' '.join(parts) + (' euro' if n == 1 else ' euros')

    # Les centimes sont toujours écrits explicitement, y compris à zéro.
    if centimes > 0:
        result += ' et ' + deux_chiffres(centimes) + ' centime' + ('' if centimes == 1 else 's')
    else:
        result += ' et zéro centime'
    return result


# ============================================================
# === LECTURE DU DET EXCEL ===================================
# ============================================================

def lire_det(chemin_excel: str) -> dict:
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl manquant. Lance : pip install openpyxl")

    wb = openpyxl.load_workbook(chemin_excel, data_only=True)

    if 'Infos étude' not in wb.sheetnames:
        raise ValueError(f"Feuille 'Infos étude' introuvable dans {chemin_excel}. "
                         f"Feuilles disponibles : {wb.sheetnames}")

    ws = wb['Infos étude']

    def val(cell_ref):
        v = ws[cell_ref].value
        return str(v).strip() if v is not None else ""

    def num(cell_ref):
        v = ws[cell_ref].value
        try:
            return float(v) if v else 0.0
        except (TypeError, ValueError):
            return 0.0

    data = {
        "numero_etude":       val("B4"),
        "nom_etude":          val("B5"),
        "type_etude":         val("B6"),
        "taille_entreprise":  val("B7"),
        "nom_entreprise":     val("B8"),
        "adresse_entreprise": val("B9"),
        "cp_ville":           val("B10"),
        "nom_client":         val("B11"),
        "chef_projet_1":      val("B12"),
        "chef_projet_2":      val("B13"),
        "duree_semaines":     val("B14"),
        "nb_qs_phase1":       val("B15"),
        "nb_qs_phase2":       val("B16"),
        "statut_etude":       val("B17"),
    }

    phases = {
        "preparation":   {"col": "B"},
        "terrain_qs":    {"col": "C"},
        "terrain_esd":   {"col": "D"},
        "analyse_bp":    {"col": "E"},
        "analyse_bench": {"col": "F"},
        "analyse_qs":    {"col": "G"},
        "analyse_esd":   {"col": "H"},
        "analyse_reco":  {"col": "I"},
    }

    for phase_key, p in phases.items():
        c = p["col"]
        phases[phase_key] = {
            "nb_responsables": num(f"{c}20"),
            "nb_jeh":          num(f"{c}21"),
            "jeh_montant":     num(f"{c}22"),
            "total_jeh_verse": num(f"{c}23"),
            "marge_jeh":       num(f"{c}24"),
            "jeh_facture":     num(f"{c}25"),
            "fact_total_jeh":  num(f"{c}26"),
        }

    data["phases"] = phases
    data["total_jeh"]          = num("J26")
    data["frais_gestion"]      = num("B29")
    data["autres_frais"]       = num("B30") if ws["B30"].value else 0.0
    data["marge_brute_pct"]    = num("B32")
    data["marge_nette_pct"]    = num("B33")
    data["pct_acompte"]        = num("B35")
    data["acompte_ht"]         = num("G9")
    data["acompte_ttc"]        = num("H9")
    data["solde_ht"]           = num("G13")
    data["solde_ttc"]          = num("H13")

    if 'Devis à exporter vers la CE' in wb.sheetnames:
        ws_devis = wb['Devis à exporter vers la CE']
        data["total_ht"]  = ws_devis["E32"].value or (data["acompte_ht"] + data["solde_ht"])
        data["total_ttc"] = ws_devis["E36"].value or (data["acompte_ttc"] + data["solde_ttc"])
    else:
        data["total_ht"]  = data["acompte_ht"] + data["solde_ht"]
        data["total_ttc"] = data["acompte_ttc"] + data["solde_ttc"]

    phases_actives = [k for k, v in phases.items() if v["nb_responsables"] > 0 or v["fact_total_jeh"] > 0]
    data["phases_actives"] = phases_actives
    data["nb_intervenants_total"] = int(num("J20"))

    return data


# ============================================================
# === FORMATAGE ===============================================
# ============================================================

def formater_montant(valeur: float) -> str:
    if valeur == int(valeur):
        return f"{int(valeur):,}".replace(",", " ")
    return f"{valeur:,.2f}".replace(",", " ").replace(".", ",")


def formater_montant_decimal(valeur: float) -> str:
    return f"{valeur:,.2f}".replace(",", " ").replace(".", ",")


def nombre_en_aine(n: int) -> str:
    mapping = {10: "dizaine", 20: "vingtaine", 30: "trentaine", 40: "quarantaine",
               50: "cinquantaine", 60: "soixantaine", 100: "centaine"}
    rounded = round(n / 10) * 10
    return mapping.get(rounded, str(rounded))


def formater_ou_multi(values: list, suffix: str = "") -> str:
    if not values:
        return "[À COMPLÉTER]" + suffix
    vals = [str(v) for v in values]
    if len(vals) == 1:
        return vals[0] + suffix
    elif len(vals) == 2:
        return f"{vals[0]} ou {vals[1]}" + suffix
    else:
        return f"{', '.join(vals[:-1])} ou {vals[-1]}" + suffix


# ============================================================
# === BRIEF IA ================================================
# ============================================================

def construire_brief_depuis_resume(resume: str, dets: list) -> str:
    brief = f"""RÉSUMÉ DU PROJET (source principale) :
{resume}
"""
    for i, data in enumerate(dets, start=1):
        phases = data["phases"]
        desc_phases = []
        if phases["terrain_esd"]["fact_total_jeh"] > 0:
            desc_phases.append(
                f"Phase terrain ESD : {int(phases['terrain_esd']['nb_responsables'])} intervenant(s), "
                f"{int(phases['terrain_esd']['nb_jeh'])} JEH à {int(phases['terrain_esd']['jeh_montant'])}€"
            )
        if phases["terrain_qs"]["fact_total_jeh"] > 0:
            desc_phases.append(
                f"Phase terrain QS : {int(phases['terrain_qs']['nb_responsables'])} intervenant(s), "
                f"{int(phases['terrain_qs']['nb_jeh'])} JEH à {int(phases['terrain_qs']['jeh_montant'])}€"
            )
        if phases["analyse_esd"]["fact_total_jeh"] > 0:
            desc_phases.append(f"Phase analyse ESD : {int(phases['analyse_esd']['nb_responsables'])} intervenant(s)")
        if phases["analyse_qs"]["fact_total_jeh"] > 0:
            desc_phases.append(f"Phase analyse QS : {int(phases['analyse_qs']['nb_responsables'])} intervenant(s)")
        if phases["analyse_reco"]["fact_total_jeh"] > 0:
            desc_phases.append(
                f"Phase conclusions & recommandations : {int(phases['analyse_reco']['nb_responsables'])} intervenant(s)"
            )

        brief += f"""
DONNÉES DET — OPTION {i} :
Nom de l'étude        : {data['nom_etude'] or 'Non renseigné'}
Type d'étude          : {data['type_etude'] or 'Non renseigné'}
Entreprise cliente    : {data['nom_entreprise'] or 'Non renseigné'}
Taille entreprise     : {data['taille_entreprise'] or 'Non renseigné'}
Contact client        : {data['nom_client'] or 'Non renseigné'}
Chef de projet        : {data['chef_projet_1'] or 'Non renseigné'}
Chef de projet 2      : {data['chef_projet_2'] or 'Aucun'}
Durée                 : {data['duree_semaines']} semaines
Nb questionnaires 1   : {data['nb_qs_phase1'] or 'Non renseigné'}
Nb questionnaires 2   : {data['nb_qs_phase2'] or 'Non renseigné'}
Structure mission     : {chr(10).join(desc_phases) if desc_phases else 'Non renseigné'}
Nb total intervenants : {data['nb_intervenants_total']}
Total JEH             : {int(data['total_jeh'])}
Montant HT            : {formater_montant(data['total_ht'])} €
Montant TTC           : {formater_montant(data['total_ttc'])} €
Acompte HT            : {formater_montant(data['acompte_ht'])} €
Acompte TTC           : {formater_montant(data['acompte_ttc'])} €
Solde HT              : {formater_montant(data['solde_ht'])} €
Solde TTC             : {formater_montant(data['solde_ttc'])} €
Marge brute           : {round(data['marge_brute_pct'] * 100, 1)}%
"""

    if not dets:
        brief += "\n(Aucun DET fourni. Les informations financières ne sont pas connues.)\n"

    return brief


# ============================================================
# === BALISES IA =============================================
# ============================================================

BALISES_IA = {
    "titre_etude":          "Titre court de l'étude (ex: Étude de marché — Secteur X).",
    "contexte_marche":      "Paragraphe de contexte pour la slide 'CONTEXTE ET ENJEUX'. CONTRAINTE DE LONGUEUR STRICTE : maximum 4 phrases et 75 mots au total (statistiques incluses). Structure : (1) commence OBLIGATOIREMENT par 'Le client' et explique brièvement ce qu'il veut et dans quel contexte de marché. (2) Intègre EXACTEMENT deux statistiques de marché chiffrées, chacune suivie de sa source entre parenthèses — chiffres réalistes et généraux sur le secteur (taille du marché, croissance annuelle…) avec sources plausibles (ex: '(Xerfi, 2024)', '(INSEE, 2023)', '(Statista, 2024)'). Paragraphe fluide et professionnel, sans titre, sans puces. Ne dépasse JAMAIS 75 mots.",
    "objet_final_etude":    "FRAGMENT DÉVELOPPÉ de 70 à 90 mots (3 à 4 lignes). Ta réponse s'affiche JUSTE APRÈS les mots 'le client souhaite' déjà présents dans la phrase. Commence DIRECTEMENT par un verbe à l'infinitif, NE RÉPÈTE JAMAIS 'le client souhaite'. Développe en détail : (1) l'objectif concret visé, (2) les résultats attendus de l'étude, (3) les décisions stratégiques que ces résultats permettront de prendre, et (4) la valeur apportée à l'entreprise. Enchaîne plusieurs propositions reliées par des virgules et 'afin de'/'ainsi que'. Ex: 'évaluer la satisfaction de ses abonnés et identifier les leviers de rétention, afin de comprendre les facteurs de fidélisation, d'orienter ses décisions stratégiques et d'optimiser durablement son offre pour renforcer sa position sur le marché'. Sans majuscule initiale, sans point final.",
    "obj_etudede":          "FRAGMENT de phrase s'insérant après 'afin de connaître'. Exemples de rendu attendu : 'le marché de la restauration rapide en Île-de-France et les attentes des consommateurs' ou 'les habitudes d'achat des ménages normands concernant les produits bio'. NE commence PAS ta réponse par 'afin de connaître'. Commence directement par l'article (le, la, les, l') ou le nom. Sans majuscule, sans point final.",
    "objectif_principal":   "1 phrase impactante résumant l'objectif principal. 20 mots max.",
    "sous_objun":           "35 à 45 mots, phrase(s) complète(s). PRIORITÉ ABSOLUE : si le brief fournit des sous-objectifs explicites (liste numérotée, ex: '1. ... 2. ... 3. ...'), REFORMULE FIDÈLEMENT le 1er sous-objectif listé — reprends son thème et son contenu exact, ne change pas le sujet, développe-le juste en phrase(s) complète(s). Si aucun sous-objectif n'est fourni dans le brief, invente-en un cohérent avec l'étude (ce qu'on cherche à mesurer, comprendre ou analyser).",
    "sous_objdeux":         "35 à 45 mots, phrase(s) complète(s). PRIORITÉ ABSOLUE : si le brief fournit des sous-objectifs explicites numérotés, REFORMULE FIDÈLEMENT le 2ème sous-objectif listé — reprends son thème exact sans le changer. Si aucun sous-objectif n'est fourni, invente-en un distinct du 1er et du 3ème (comportements d'achat, perception, attentes, profil de la cible, etc.).",
    "sous_objtrois":        "35 à 45 mots, phrase(s) complète(s). PRIORITÉ ABSOLUE : si le brief fournit des sous-objectifs explicites numérotés, REFORMULE FIDÈLEMENT le 3ème sous-objectif listé — reprends son thème exact sans le changer. Si aucun sous-objectif n'est fourni, invente-en un distinct des deux premiers (recommandations, positionnement, opportunités, fidélisation, etc.).",
    "info_questionnaire":   "FRAGMENT. Complète 'Une [X]aine de données à majorité quantitative/qualitative'. Déduis le type depuis le brief. Ex: 'Une vingtaine de données à majorité quantitative'.",
    "of_questionnaire":     "FRAGMENT. Complète 'Une [X]aine de questions à majorité fermées/ouvertes'. Déduis depuis le brief. Ex: 'Une vingtaine de questions à majorité fermées'.",
    "tps_questionnaire":    "FRAGMENT. Complète 'durée comprise entre...'. Ex: '5 et 7 minutes'.",
    "nb_question":          "FRAGMENT. Complète 'au maximum de...'. Ex: '15 questions dont 12 fermées'.",
    "obj_etudetr":          "FRAGMENT. Complète 'de préciser...'. Ex: 'de préciser les attentes des consommateurs quant à l'offre de service'. Sans majuscule, sans point final.",
    "secteur":              "FRAGMENT s'insérant après 'une analyse précise du secteur'. Donne SOIT un adjectif de secteur ('immobilier', 'bancaire', 'automobile', 'cosmétique'), SOIT un complément introduit par de la/du/de l' ('de la restauration rapide', 'de la grande distribution'). Sans majuscule, sans point final. Ne répète PAS le mot 'secteur'.",
    "produit_service":      "FRAGMENT. NOM PROPRE UNIQUEMENT du produit ou service étudié. Exemples : 'Proper', 'iPhone', 'Netflix'. N'ajoute AUCUN mot générique ('solution', 'produit', 'service', article, etc.). Juste le nom propre.",
    "marche_pot":           "Rédige un paragraphe fluide et professionnel de 3-4 phrases décrivant l'étude documentaire du marché. Explique qu'elle permet de définir le marché potentiel et de réaliser un diagnostic de l'offre existante dans le secteur du client. Précise qu'elle repose sur un tour d'horizon des dernières informations disponibles, et annonce une analyse de l'état et des tendances du marché ainsi que de sa situation géographique. Adapte au secteur et à la zone du brief. AUCUN crochet, AUCune parenthèse d'instruction, AUCun placeholder dans ta réponse.",
    "est_client":           "FRAGMENT. Estimation du marché potentiel (chiffre ou fourchette).",
    "concu_direct":         "FRAGMENT. 2 ou 3 concurrents directs majeurs opérant dans la zone géographique et le secteur du projet. OBLIGATOIRE : cite des noms réels basés sur tes connaissances du secteur même si non mentionnés dans le brief. N'écris JAMAIS '[À COMPLÉTER]' pour cette balise. Ex: 'Darty, Boulanger, Fnac'.",
    "concu_indi":           "FRAGMENT. 2 ou 3 concurrents indirects opérant dans la zone et le secteur. OBLIGATOIRE : cite des noms réels basés sur tes connaissances même si non mentionnés dans le brief. N'écris JAMAIS '[À COMPLÉTER]' pour cette balise. Ex: 'Amazon, Cdiscount, LeBonCoin'.",
    "analyse_donnee":       "Phrase nominale autonome (encadré affiché parmi les objectifs) : 'Une analyse précise de [aspect spécifique à l'étude]'. Sois SPÉCIFIQUE au projet du brief, pas générique. Ex: 'Une analyse précise des comportements d'achat des consommateurs havrais'. Commence par une majuscule, sans point final.",
    "cible":                "FRAGMENT. Groupe nominal décrivant la cible visée (avec article). Ex: 'les ménages de la région Normandie'.",
    "zone_action":          "FRAGMENT. Zone géographique de l'étude. Ex: 'la région Île-de-France'.",
    "zone_passation":       "FRAGMENT. Zone de passation des questionnaires.",
    "lieu_client":          "FRAGMENT. Lieu où se trouve la clientèle cible.",
    "selection_region":     "1 phrase justifiant le choix de la zone géographique.",
    "justif_cible":         "1 phrase justifiant la pertinence de la cible choisie. Compléter la phrase : Le listing veillera à cibler",
    "cible_listing":        "FRAGMENT. Description courte de la cible pour le listing (ex: 'des responsables marketing de PME franciliennes'). UNE SEULE description sans format Option 1/Option 2.",
    "cible_ville":          "FRAGMENT. 3 exemples de villes représentatives de la zone d'action de l'étude, séparées par des virgules. Ex: 'Rouen, Caen et Le Havre'.",
    "ressource_com":        "15 mots max : Ressources de communication recommandées (journaux, magazines spécialisés, réseaux sociaux...).",
    "action_com":           "30 mots environ : Décris les principales actions de communication concrètes à mettre en place pour atteindre les objectifs de visibilité et de notoriété (ex: présence sur réseaux sociaux, partenariats médias, campagnes ciblées, relations presse...).",
    "obj_etude":            "FRAGMENT. Complète 'Afin de...'. Verbe infinitif.",
    "satis_globale":        "FRAGMENT. Complète 'La satisfaction globale [ta réponse]'. En 1 phrase courte, explique la satisfaction globale de la cible concernant la solution/produit/service du client. Ex: 'vis-à-vis de la solution de livraison proposée, notamment en termes de délais et de fiabilité'. Sans majuscule initiale.",
    "besoinetatt":          "FRAGMENT. Complète 'Les besoins et les attentes [ta réponse]'. En 1 phrase courte, décris les besoins et attentes de la cible vis-à-vis du produit/service. Ex: 'des consommateurs concernant l'offre en termes de praticité, prix et qualité perçue'. Sans majuscule initiale.",
    "avis_etude":           "FRAGMENT. Complète '...quant à...'. Reformule les objectifs de l'étude sous forme d'avis/opinions à recueillir. Ex: 'l'offre de service proposée et son positionnement sur le marché'.",
    "besoin_client":        "Rédige 2 à 3 phrases fluides et professionnelles expliquant que le guide d'entretien permettra d'orienter la discussion sur la satisfaction, les attentes et les besoins des clients ou consommateurs potentiels concernant le produit ou service du client. Développe avec les thèmes spécifiques abordés (qualité, prix, usage, amélioration souhaitée). Tu peux reprendre et reformuler les sous-objectifs. AUCUN crochet, parenthèse d'instruction ni placeholder dans ta réponse.",
    "besoinclient":         "1 à 2 phrases décrivant le besoin principal du client pour cette étude. Explique clairement ce qu'il cherche à accomplir, résoudre ou améliorer grâce aux résultats de l'étude.",
    "obj_final_my":         "FRAGMENT. Complète 'mieux appréhender...'. Ex: 'mieux appréhender la qualité de l'accueil en point de vente'. Sans majuscule, sans point final.",
    "obj_myst":             "Phrase COMPLÈTE et autonome (commence par une majuscule, finit par un point). Décris l'objectif majeur des visites mystères pour ce projet. Ex: 'L'objectif principal de ces visites mystères est d'évaluer la qualité de l'accueil et du conseil en point de vente afin d'identifier les axes d'amélioration prioritaires.'",
    "produitouserv":        "FRAGMENT. Écris UNIQUEMENT le mot 'produit' ou le mot 'service' (selon ce dont il s'agit dans le brief). Rien d'autre, aucun autre mot.",
    "poste_concer":         "FRAGMENT. Type de contact/poste visé au sein des entreprises du listing. Complète la phrase '...contacts privilégiés au sein de ces entreprises...' Ex: '(directeurs marketing, responsables RH, gérants)'.",
    "listing":              "Complète la phrase qui commence par 'Neoma Conseil s'engage à réaliser un listing quantitatif et qualitatif d'entreprises'. Ta réponse suit directement ce début et précise le type d'entreprises ciblées (adapté au secteur et au besoin du client), susceptibles d'être intéressées par le produit/service. Ex: 'spécialisées dans l'immobilier résidentiel, susceptibles de recourir à une solution de création de vidéos'. Sans placeholder ni '(à adapter)'.",
    "livrable":             "2 éléments MAXIMUM constituant le rapport final, séparés par des sauts de ligne (\\n). SANS tiret ni puce ni numérotation. Chaque élément est une courte phrase nominale. Ex: 'Analyse documentaire du secteur et des tendances\\nRésultats de l'enquête et recommandations stratégiques'. Adapte au type d'étude du brief.",
    "validite_date":        "Date de validité (3 mois après aujourd'hui). Format : JJ mois AAAA.",
    "quota_cat1_opt1":      "FRAGMENT. Quota catégorie 1 option 1. Juste le chiffre.",
    "quota_cat1_opt2":      "FRAGMENT. Quota catégorie 1 option 2. Juste le chiffre.",
    "quota_cat1_opt3":      "FRAGMENT. Quota catégorie 1 option 3. Juste le chiffre.",
    "quota_cat2_opt1":      "FRAGMENT. Quota catégorie 2 option 1. Juste le chiffre.",
    "quota_cat2_opt2":      "FRAGMENT. Quota catégorie 2 option 2. Juste le chiffre.",
    "quota_cat2_opt3":      "FRAGMENT. Quota catégorie 2 option 3. Juste le chiffre.",
    "quota_cat3_opt1":      "FRAGMENT. Quota catégorie 3 option 1. Juste le chiffre.",
    "quota_cat3_opt2":      "FRAGMENT. Quota catégorie 3 option 2. Juste le chiffre.",
    "quota_cat3_opt3":      "FRAGMENT. Quota catégorie 3 option 3. Juste le chiffre.",
    "quota_cat4_opt1":      "FRAGMENT. Quota catégorie 4 option 1. Juste le chiffre.",
    "quota_cat4_opt2":      "FRAGMENT. Quota catégorie 4 option 2. Juste le chiffre.",
    "quota_cat4_opt3":      "FRAGMENT. Quota catégorie 4 option 3. Juste le chiffre.",
    "pct_cat1":             "FRAGMENT. Pourcentage catégorie 1, sans %. Ex: 25",
    "pct_cat2":             "FRAGMENT. Pourcentage catégorie 2, sans %. Ex: 25",
    "pct_cat3":             "FRAGMENT. Pourcentage catégorie 3, sans %. Ex: 25",
    "pct_cat4":             "FRAGMENT. Pourcentage catégorie 4, sans %. Ex: 25",
    "tel_client":           "FRAGMENT. Numéro de téléphone du client format français. Si absent du brief: [À COMPLÉTER].",
    "mail_client":          "Adresse e-mail du client. Si absente du brief: [À COMPLÉTER].",
    "tel_chef_projet":      "FRAGMENT. Numéro de téléphone du chef de projet format français.",
    "nb_listing":           "FRAGMENT. Nombre de contacts dans le listing. Juste le chiffre.",
    "nb_repondants":        "Nombre de répondants cibles pour l'étude. Juste le chiffre.",
    "questioprinter":       "FRAGMENT. Nombre de questionnaires que chaque intervenant devra administrer. Juste le chiffre (ex: '8'). Si pas de DET, estime en divisant le nombre total de questionnaires par le nombre d'intervenants.",
    "nb_ville":             "FRAGMENT. Nombre de villes ou zones géographiques sélectionnées pour la passation. Juste le chiffre, sans texte autour (ex: '3'). Déduis-le de la zone géographique et du volume de questionnaires du brief.",
    "nom_entreprise":       "Nom de l'entreprise cliente.",
    "nom_representant":     "Prénom et nom du représentant client.",
    "etablisseur_je":       "Prénom et nom du chef de projet Neoma Conseil.",
    "validateur_je":        "Prénom et nom du validateur Neoma Conseil.",
    "intervenant":          "Prénom(s) des intervenants.",
    "nb_intervenants":      "Nombre total d'intervenants. Juste le chiffre.",
    "qualiteje":            "2 phrases décrivant les compétences et appétences du chef de projet en lien direct avec cette étude. Ex: 'Passionné par le secteur de la restauration, il dispose d'une expérience en étude de marché alimentaire. Sa maîtrise des outils d'analyse quantitative garantit des résultats fiables.'",
    "nb_questionnaire":     "Nombre de questionnaires. Ex: 150.",
    "nb_questionnaire_min": "Nombre minimum de questionnaires.",
    "nb_questionnaire_max": "Nombre maximum de questionnaires.",
    # ── Option 1 (noms alignés sur le template) ──
    "montant_acompte_htun":   "Montant acompte HT option 1, format '1 250,00'. Si aucune donnée financière : [À COMPLÉTER].",
    "montant_acompte_ttcun":  "Montant acompte TTC option 1, format '1 500,00'. Si aucune donnée : [À COMPLÉTER].",
    "montant_solde_htun":     "Montant solde HT option 1, format '1 250,00'. Si aucune donnée : [À COMPLÉTER].",
    "montant_solde_ttcun":    "Montant solde TTC option 1, format '1 500,00'. Si aucune donnée : [À COMPLÉTER].",
    "Duree_etudeun":          "Durée de l'étude option 1. Ex: '8 semaines'.",
    # ── Option 2 ──
    "montant_acompte_htde":   "Montant acompte HT option 2, format '1 250,00'. Si aucune donnée : [À COMPLÉTER].",
    "montant_acompte_ttcde":  "Montant acompte TTC option 2, format '1 500,00'. Si aucune donnée : [À COMPLÉTER].",
    "montant_solde_htde":     "Montant solde HT option 2, format '1 250,00'. Si aucune donnée : [À COMPLÉTER].",
    "montant_solde_ttcde":    "Montant solde TTC option 2, format '1 500,00'. Si aucune donnée : [À COMPLÉTER].",
    "Duree_etudede":          "Durée de l'étude option 2. Ex: '10 semaines'.",
    "nb_questionnairede":     "Nombre de questionnaires option 2. Juste le chiffre.",
    # ── Option 3 ──
    "montant_acompte_httrois":  "Montant acompte HT option 3, format '1 250,00'. Si aucune donnée : [À COMPLÉTER].",
    "montant_acompte_ttctrois": "Montant acompte TTC option 3, format '1 500,00'. Si aucune donnée : [À COMPLÉTER].",
    "montant_solde_httrois":    "Montant solde HT option 3, format '1 250,00'. Si aucune donnée : [À COMPLÉTER].",
    "montant_solde_ttctrois":   "Montant solde TTC option 3, format '1 500,00'. Si aucune donnée : [À COMPLÉTER].",
    "duree_etudetrois":         "Durée de l'étude option 3. Ex: '12 semaines'.",
    "nb_questionnairetr":       "Nombre de questionnaires option 3. Juste le chiffre.",
}

# Balises montants dont la version "en lettres" ({{...ltr}}) est calculée automatiquement
MONTANTS_AVEC_LTR = [
    "montant_acompte_htun",  "montant_acompte_ttcun",  "montant_solde_htun",  "montant_solde_ttcun",
    "montant_acompte_htde",  "montant_acompte_ttcde",  "montant_solde_htde",  "montant_solde_ttcde",
    "montant_acompte_httrois", "montant_acompte_ttctrois", "montant_solde_httrois", "montant_solde_ttctrois",
]


# Amorces de phrase présentes dans le template : si l'IA les répète en tête
# de sa réponse, on les retire pour éviter les doublons ('le client souhaite
# le client souhaite …').
AMORCES_REDONDANTES = {
    "{{objet_final_etude}}": "le client souhaite",
    "{{obj_etudede}}":       "afin de connaître",
    "{{obj_etudetr}}":       "de préciser",
    "{{obj_final_my}}":      "dans le but de",
    "{{obj_etude}}":         "afin de",
    "{{satis_globale}}":     "la satisfaction globale",
    "{{besoinetatt}}":       "les besoins et les attentes",
    "{{avis_etude}}":        "quant à",
}


def nettoyer_fragments(replacements: dict):
    """Retire les amorces de phrase que l'IA aurait répétées, et corrige la
    ponctuation résiduelle en tête de fragment."""
    for tag, amorce in AMORCES_REDONDANTES.items():
        val = replacements.get(tag)
        if not val:
            continue
        # Retirer l'amorce répétée en tête (insensible à la casse/accents simples)
        low = val.lstrip()
        if low.lower().startswith(amorce.lower()):
            val = low[len(amorce):].lstrip(" ,':")
        # Nettoyer une éventuelle majuscule initiale parasite sur un fragment
        replacements[tag] = val.strip()


def parser_montant_fr(val: str):
    """Parse un montant français ('1 250,00', '1250', '2 500 €') → float, ou None."""
    if not val:
        return None
    cleaned = re.sub(r'[^\d,\.]', '', str(val)).replace(',', '.')
    # Gérer '1.250.00' résiduel : garder le dernier point comme séparateur décimal
    parts = cleaned.split('.')
    if len(parts) > 2:
        cleaned = ''.join(parts[:-1]) + '.' + parts[-1]
    try:
        return float(cleaned) if cleaned else None
    except ValueError:
        return None


def completer_montants_en_lettres(replacements: dict):
    """
    Normalise le format numérique des montants (toujours 2 décimales, ex:
    '800,00' même pour un montant rond) et complète les balises {{...ltr}}
    manquantes depuis ces montants numériques.
    """
    for base in MONTANTS_AVEC_LTR:
        tag, ltr_tag = "{{" + base + "}}", "{{" + base + "ltr}}"
        if tag not in replacements:
            continue
        val = replacements[tag]
        if "[À COMPLÉTER]" in str(val):
            replacements.setdefault(ltr_tag, "[À COMPLÉTER]")
            continue
        num = parser_montant_fr(val)
        if num is None or num <= 0:
            continue
        # Garantir 2 décimales même si l'IA a répondu un montant rond ('800' -> '800,00')
        replacements[tag] = formater_montant_decimal(num)
        if ltr_tag not in replacements:
            replacements[ltr_tag] = montant_en_lettres(num)


def appeler_ia(brief: str, api_key: str, model: str = GROQ_MODEL, contextes: dict = None) -> dict:
    balises = get_balises_effectives()
    contextes = contextes or {}
    lines = []
    for k, v in balises.items():
        tag  = "{{" + k + "}}"
        line = f"{tag}: {v}"
        # Pour les fragments, donner le contexte d'insertion réel du template
        ctx = contextes.get(tag)
        if ctx and "FRAGMENT" in v:
            line += f" [Contexte: {ctx}]"
        lines.append(line)

    prompt = f"""Tu es un consultant expert chez Neoma Conseil, spécialisé dans les propositions commerciales.

À partir du brief fourni, génère le contenu pour chaque balise.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
RÈGLES ABSOLUES — RESPECTE-LES TOUTES :

1. FORMAT STRICT : {{{{nom_balise}}}}: Contenu
   → Une balise par ligne. Aucune ligne supplémentaire. Aucun markdown. Aucune explication.

2. FRAGMENT = fragment de phrase : sans majuscule initiale, sans point final.

3. RÉPONDS EXCLUSIVEMENT EN FRANÇAIS PROFESSIONNEL.

4. ⚠️ ANTI-HALLUCINATION — RÈGLE LA PLUS IMPORTANTE :
   Si l'information n'est PAS présente dans le brief (résumé ou DET), tu dois écrire
   EXACTEMENT "[À COMPLÉTER]" comme valeur — N'INVENTE JAMAIS une information absente.
   Exceptions : concu_direct et concu_indi → TOUJOURS citer des noms réels basés sur tes
   connaissances du secteur, même si absent du brief.

5. EXTRACTION RIGOUREUSE : Lis intégralement le résumé avant de répondre. Toute
   information présente dans le brief DOIT être utilisée (ne mets pas [À COMPLÉTER]
   si l'information est dans le brief, même formulée différemment).

6. ⚠️ INTERDIT — N'UTILISE JAMAIS le format "Option 1 : ... / Option 2 : ..." dans
   aucune balise. Les balises sont déjà différenciées par option dans le template.
   Chaque balise doit avoir UNE SEULE valeur directement utilisable.

7. OPTIONS 2 ET 3 : si le brief ne donne qu'une seule option financière, propose des
   variantes cohérentes pour opt2/opt3 (volumes légèrement différents). Si aucune donnée
   financière n'est disponible → [À COMPLÉTER].
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

BALISES À REMPLIR :
{chr(10).join(lines)}

BRIEF :
{brief}
"""

    client = OpenAI(api_key=api_key, base_url="https://api.groq.com/openai/v1")
    resp = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.1,
        max_tokens=4000,
    )
    raw = resp.choices[0].message.content

    replacements = {}
    for key in balises:
        tag = "{{" + key + "}}"
        pattern = re.escape(tag) + r"\s*:\s*(.+?)(?=\n\s*\{\{|\Z)"
        m = re.search(pattern, raw, re.DOTALL)
        if m:
            val = re.sub(r'\*+', '', m.group(1).strip())
            val = re.sub(r'\s+', ' ', val)
            replacements[tag] = val
        else:
            pattern2 = r'(?:^|\n)\s*' + re.escape(key) + r'\s*:\s*(.+)'
            m2 = re.search(pattern2, raw)
            if m2:
                replacements[tag] = re.sub(r'\*+', '', m2.group(1).strip())

    return replacements


# ============================================================
# === REMPLACEMENTS DIRECTS DEPUIS DET =======================
# ============================================================

def construire_replacements_directs(data: dict, suffix: str = "") -> dict:
    r = {}
    s = suffix

    if not suffix:
        if data["nom_entreprise"]:
            r["{{nom_entreprise}}"]         = data["nom_entreprise"]
        if data["nom_etude"]:
            r["{{titre_etude}}"]            = data["nom_etude"]
        if data["nom_client"]:
            r["{{nom_representant}}"]       = data["nom_client"]
        if data["chef_projet_1"]:
            r["{{etablisseur_je}}"]         = data["chef_projet_1"]
        if data["chef_projet_2"]:
            r["{{validateur_je}}"]          = data["chef_projet_2"]
        if data["nb_intervenants_total"] > 0:
            nb = str(data["nb_intervenants_total"])
            r["{{nb_intervenants}}"]        = nb
            r["{{nb_intervenants_total}}"]  = nb
        if data["nb_qs_phase1"]:
            r["{{nb_questionnaire}}"]       = str(data["nb_qs_phase1"])
            r["{{nb_questionnaire_min}}"]   = str(data["nb_qs_phase1"])
            r["{{nb_repondants}}"]          = str(data["nb_qs_phase1"])
            r["{{of_questionnaire}}"]       = str(data["nb_qs_phase1"])
        if data["nb_qs_phase2"]:
            r["{{nb_questionnaire_max}}"]   = str(data["nb_qs_phase2"])
        if data["nb_intervenants_total"] > 0:
            r["{{intervenant}}"] = str(data["nb_intervenants_total"]) + " intervenants"
        elif data["chef_projet_1"]:
            intervenants = [data["chef_projet_1"].split()[0]]
            if data["chef_projet_2"]:
                intervenants.append(data["chef_projet_2"].split()[0])
            r["{{intervenant}}"] = " et ".join(intervenants)
        if data["type_etude"]:
            r["{{type_rapport}}"]           = data["type_etude"].lower()

    SUFFIXE_TAG = {"": "un", "_opt2": "de", "_opt3": "trois"}
    st = SUFFIXE_TAG.get(s, "un")

    if data["duree_semaines"]:
        duree_val = f"{data['duree_semaines']} semaines"
        r[f"{{{{Duree_etude{st}}}}}"] = duree_val
        r[f"{{{{duree_etude{st}}}}}"] = duree_val
    if data["nb_qs_phase1"]:
        r[f"{{{{nb_questionnaire{s}}}}}"]     = str(data["nb_qs_phase1"])
        r[f"{{{{nb_questionnaire_min{s}}}}}"] = str(data["nb_qs_phase1"])
        r[f"{{{{nb_questionnairetr{s}}}}}"]   = str(data["nb_qs_phase1"])
    if data["nb_qs_phase2"]:
        r[f"{{{{nb_questionnaire_max{s}}}}}"] = str(data["nb_qs_phase2"])
    if data["acompte_ht"] > 0:
        r[f"{{{{montant_acompte_ht{st}}}}}"]     = formater_montant_decimal(data["acompte_ht"])
        r[f"{{{{montant_acompte_ht{st}ltr}}}}"]  = montant_en_lettres(data["acompte_ht"])
    if data["acompte_ttc"] > 0:
        r[f"{{{{montant_acompte_ttc{st}}}}}"]    = formater_montant_decimal(data["acompte_ttc"])
        r[f"{{{{montant_acompte_ttc{st}ltr}}}}"] = montant_en_lettres(data["acompte_ttc"])
    if data["solde_ht"] > 0:
        r[f"{{{{montant_solde_ht{st}}}}}"]       = formater_montant_decimal(data["solde_ht"])
        r[f"{{{{montant_solde_ht{st}ltr}}}}"]    = montant_en_lettres(data["solde_ht"])
    if data["solde_ttc"] > 0:
        r[f"{{{{montant_solde_ttc{st}}}}}"]      = formater_montant_decimal(data["solde_ttc"])
        r[f"{{{{montant_solde_ttc{st}ltr}}}}"]   = montant_en_lettres(data["solde_ttc"])
    if suffix == "_opt2" and data["nb_qs_phase2"]:
        r["{{nb_questionnairede}}"] = str(data["nb_qs_phase2"])

    return r


# ============================================================
# === MOTEUR DE REMPLACEMENT XML =============================
# ============================================================

def apply_bold_to_rpr(rpr: str) -> str:
    if 'b="' in rpr or "b='" in rpr:
        rpr = re.sub(r'\bb="0"', 'b="1"', rpr)
        rpr = re.sub(r"b='0'", "b='1'", rpr)
    else:
        rpr = rpr.replace('<a:rPr', '<a:rPr b="1"', 1)
    return rpr


def process_paragraph(para_xml, replacements):
    run_re = re.compile(r'<a:r>.*?</a:r>', re.DOTALL)
    runs = list(run_re.finditer(para_xml))
    if not runs:
        return para_xml
    run_texts = []
    for r in runs:
        m = re.search(r'<a:t[^>]*>(.*?)</a:t>', r.group(0), re.DOTALL)
        run_texts.append(m.group(1) if m else '')
    concat = ''.join(run_texts)
    if not any(tag in concat for tag in replacements):
        return para_xml
    has_fragmented = any(
        tag in concat and not any(tag in rt for rt in run_texts)
        for tag in replacements
    )
    if has_fragmented:
        rpr_m = re.search(r'<a:rPr.*?</a:rPr>', runs[0].group(0), re.DOTALL)
        rpr = rpr_m.group(0) if rpr_m else ''
        new_text = concat
        for tag, val in replacements.items():
            new_text = new_text.replace(tag, val)
        if any(tag in concat for tag in BOLD_TAGS):
            rpr = apply_bold_to_rpr(rpr)
        new_run = f'<a:r>{rpr}<a:t xml:space="preserve">{new_text}</a:t></a:r>'
        return para_xml[:runs[0].start()] + new_run + para_xml[runs[-1].end():]
    else:
        result = para_xml
        for tag, val in replacements.items():
            result = result.replace(tag, val)
        return result


def process_xml(content, replacements):
    para_re = re.compile(r'<a:p(?:\s[^>]*)?>.*?</a:p>', re.DOTALL)
    result = para_re.sub(lambda m: process_paragraph(m.group(0), replacements), content)
    for tag, val in replacements.items():
        result = result.replace(tag, val)
    return result


def remplacer_tout(work_dir, replacements):
    replaced = set()
    para_re = re.compile(r'<a:p(?:\s[^>]*)?>.*?</a:p>', re.DOTALL)
    for path in glob.glob(os.path.join(work_dir, '**', '*.xml'), recursive=True):
        try:
            with open(path, 'r', encoding='utf-8') as f:
                content = f.read()
        except Exception:
            continue
        new_content = process_xml(content, replacements)
        if new_content != content:
            with open(path, 'w', encoding='utf-8') as f:
                f.write(new_content)
            for pm in para_re.finditer(content):
                texts = re.findall(r'<a:t[^>]*>(.*?)</a:t>', pm.group(0), re.DOTALL)
                concat = ''.join(texts)
                for tag in replacements:
                    if tag in concat:
                        replaced.add(tag)
            for tag in replacements:
                if tag in content:
                    replaced.add(tag)
    return replaced


def colorier_texte_remplace(
    pptx_bytes: bytes,
    replacements: dict,
    template_path: str,
    hex_color: str = "00B050",
) -> tuple:
    """
    Colore précisément les runs remplacés (pas tout le paragraphe).
    Gère : text frames, cellules de tableau, et group shapes (récursif).
    Place l'exposant [N] juste avant le run remplacé.
    Retourne (pptx_bytes_coloré, tag_index) où tag_index = {numéro: "{{tag}}"}
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from lxml import etree
    import io as _io, re as _re, copy as _copy

    A_NS   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    TAG_RE = _re.compile(r'\{\{([^}]+)\}\}')

    # ── Itérateur récursif : yields (location_key, paragraph) ──────────────
    # Couvre text frames, tableaux, et group shapes imbriqués.
    def iter_paras(slide, s_i):
        def process_shape(shape, path):
            # Group shape : itérer les enfants récursivement
            if hasattr(shape, 'shapes'):
                for c_i, child in enumerate(shape.shapes):
                    yield from process_shape(child, path + (c_i,))
            # Text frame classique
            elif shape.has_text_frame:
                for p_i, para in enumerate(shape.text_frame.paragraphs):
                    yield (s_i, path, 'tf', p_i), para
            # Tableau
            elif shape.has_table:
                for r_i, row in enumerate(shape.table.rows):
                    for c_i, cell in enumerate(row.cells):
                        for p_i, para in enumerate(cell.text_frame.paragraphs):
                            yield (s_i, path, 'tbl', r_i, c_i, p_i), para

        for sh_i, shape in enumerate(slide.shapes):
            yield from process_shape(shape, (sh_i,))

    # ── Étape 1 : lire le template, indexer les paragraphes contenant des tags ──
    tagged_runs: dict = {}

    prs_tmpl = Presentation(template_path)
    for s_i, slide in enumerate(prs_tmpl.slides):
        for key, para in iter_paras(slide, s_i):
            run_info       = []
            tmpl_run_texts = []
            for r_i, run in enumerate(para.runs):
                tmpl_run_texts.append(run.text)
                raw_tags = TAG_RE.findall(run.text)
                found = list(dict.fromkeys(
                    "{{" + t + "}}" for t in raw_tags
                    if "{{" + t + "}}" in replacements
                ))
                if found:
                    run_info.append((r_i, found))

            if not run_info:
                # Tag éclaté sur plusieurs runs (fragmentation XML)
                full_text = "".join(tmpl_run_texts)
                frag_tags = list(dict.fromkeys(
                    "{{" + t + "}}" for t in TAG_RE.findall(full_text)
                    if "{{" + t + "}}" in replacements
                    and not any("{{" + t + "}}" in rt for rt in tmpl_run_texts)
                ))
                if frag_tags:
                    run_info = [(-1, frag_tags)]  # -1 = marqueur "fragmenté"

            if run_info:
                tagged_runs[key] = {
                    'run_info':   run_info,
                    'tmpl_texts': tmpl_run_texts,
                }

    # ── Étape 2 : numérotation par ordre d'apparition ──────────────────────
    tag_to_num: dict = {}
    counter = [1]

    def get_num(tag: str) -> int:
        if tag not in tag_to_num:
            tag_to_num[tag] = counter[0]
            counter[0] += 1
        return tag_to_num[tag]

    for key in sorted(tagged_runs):
        for r_i, tags in tagged_runs[key]['run_info']:
            for tag in tags:
                get_num(tag)

    # ── Helpers XML ─────────────────────────────────────────────────────────
    def creer_expose(number: int) -> etree._Element:
        r   = etree.Element(f'{{{A_NS}}}r')
        rPr = etree.SubElement(r, f'{{{A_NS}}}rPr')
        rPr.set('baseline', '30000')
        rPr.set('sz', '800')
        rPr.set('b', '0')
        fill = etree.SubElement(rPr, f'{{{A_NS}}}solidFill')
        clr  = etree.SubElement(fill, f'{{{A_NS}}}srgbClr')
        clr.set('val', hex_color)
        t    = etree.SubElement(r, f'{{{A_NS}}}t')
        t.text = f'[{number}]'
        return r

    def colorier_rpr(r_elem: etree._Element):
        rPr = r_elem.find(f'{{{A_NS}}}rPr')
        if rPr is None:
            rPr = etree.Element(f'{{{A_NS}}}rPr')
            r_elem.insert(0, rPr)
        for ft in ('solidFill', 'gradFill', 'pattFill'):
            ex = rPr.find(f'{{{A_NS}}}{ft}')
            if ex is not None:
                rPr.remove(ex)
        fill = etree.SubElement(rPr, f'{{{A_NS}}}solidFill')
        clr  = etree.SubElement(fill, f'{{{A_NS}}}srgbClr')
        clr.set('val', hex_color)

    def remplacer_run_par_split(p_elem, r_elem, text, val, num):
        idx = text.find(val) if val else -1
        if idx == -1:
            colorier_rpr(r_elem)
            r_idx = list(p_elem).index(r_elem)
            p_elem.insert(r_idx, creer_expose(num))
            return
        new_elems = []
        if idx > 0:
            r_pre = _copy.deepcopy(r_elem)
            r_pre.find(f'{{{A_NS}}}t').text = text[:idx]
            new_elems.append(r_pre)
        new_elems.append(creer_expose(num))
        r_val = _copy.deepcopy(r_elem)
        r_val.find(f'{{{A_NS}}}t').text = val
        colorier_rpr(r_val)
        new_elems.append(r_val)
        suffix = text[idx + len(val):]
        if suffix:
            r_suf = _copy.deepcopy(r_elem)
            r_suf.find(f'{{{A_NS}}}t').text = suffix
            new_elems.append(r_suf)
        r_idx = list(p_elem).index(r_elem)
        p_elem.remove(r_elem)
        for k, elem in enumerate(new_elems):
            p_elem.insert(r_idx + k, elem)

    def colorier_para(para, key):
        """Applique la coloration à un paragraphe du PPTX généré."""
        info      = tagged_runs[key]
        run_info  = info['run_info']
        tmpl_txts = info['tmpl_texts']
        gen_runs  = list(para.runs)
        p_elem    = para._p

        # ── Cas fragmenté (run_idx == -1) ──────────────────────────────────
        if run_info[0][0] == -1:
            tags_frag = run_info[0][1]
            for g_run in list(para.runs):
                run_text = g_run.text
                if not run_text:
                    continue
                entries = sorted(
                    [(run_text.find(replacements.get(tag, "")), tag)
                     for tag in tags_frag
                     if replacements.get(tag, "") and replacements[tag] in run_text],
                    key=lambda x: x[0]
                )
                if not entries:
                    if run_text.strip():
                        g_run.font.color.rgb = RGBColor(
                            int(hex_color[:2], 16),
                            int(hex_color[2:4], 16),
                            int(hex_color[4:], 16),
                        )
                    continue
                cursor    = 0
                new_elems = []
                for pos, tag in entries:
                    val = replacements[tag]
                    if pos > cursor:
                        r_pre = _copy.deepcopy(g_run._r)
                        r_pre.find(f'{{{A_NS}}}t').text = run_text[cursor:pos]
                        new_elems.append(r_pre)
                    new_elems.append(creer_expose(get_num(tag)))
                    r_val = _copy.deepcopy(g_run._r)
                    r_val.find(f'{{{A_NS}}}t').text = val
                    colorier_rpr(r_val)
                    new_elems.append(r_val)
                    cursor = pos + len(val)
                if cursor < len(run_text):
                    r_suf = _copy.deepcopy(g_run._r)
                    r_suf.find(f'{{{A_NS}}}t').text = run_text[cursor:]
                    new_elems.append(r_suf)
                r_idx = list(p_elem).index(g_run._r)
                p_elem.remove(g_run._r)
                for k, elem in enumerate(new_elems):
                    p_elem.insert(r_idx + k, elem)
            return

        # ── Cas non fragmenté : traitement run par run ──────────────────────
        for r_i, tags_in_run in run_info:
            if r_i >= len(gen_runs):
                continue
            g_run     = gen_runs[r_i]
            gen_text  = g_run.text
            tmpl_text = tmpl_txts[r_i] if r_i < len(tmpl_txts) else ""
            tmpl_sans_tags = TAG_RE.sub("", tmpl_text).strip()

            for tag in tags_in_run:
                val = replacements.get(tag, "")
                num = get_num(tag)
                if not tmpl_sans_tags:
                    g_run.font.color.rgb = RGBColor(
                        int(hex_color[:2], 16),
                        int(hex_color[2:4], 16),
                        int(hex_color[4:], 16),
                    )
                    r_idx = list(p_elem).index(g_run._r)
                    p_elem.insert(r_idx, creer_expose(num))
                else:
                    remplacer_run_par_split(p_elem, g_run._r, gen_text, val, num)
                break

    # ── Étape 3 : annoter le PPTX généré ───────────────────────────────────
    prs = Presentation(_io.BytesIO(pptx_bytes))
    for s_i, slide in enumerate(prs.slides):
        for key, para in iter_paras(slide, s_i):
            if key in tagged_runs:
                try:
                    colorier_para(para, key)
                except Exception:
                    pass  # Ne jamais planter la génération pour une erreur de coloration

    tag_index = {num: tag for tag, num in tag_to_num.items()}
    out = _io.BytesIO()
    prs.save(out)
    return out.getvalue(), tag_index


# ============================================================
# === DÉFRAGMENTATION DES BALISES ============================
# ============================================================

def defragmenter_pptx(pptx_bytes: bytes) -> bytes:
    """
    Fusionne les runs XML consécutifs lorsqu'une balise {{...}} est éclatée
    sur plusieurs runs (artefact d'export Google Slides / correcteur ortho).
    Après ce traitement, chaque balise tient dans UN SEUL run, ce qui rend
    le remplacement et la coloration fiables à 100 %.
    """
    from lxml import etree
    import io as _io

    A_NS   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    R_TAG  = f'{{{A_NS}}}r'
    T_TAG  = f'{{{A_NS}}}t'
    P_TAG  = f'{{{A_NS}}}p'
    TAG_SPAN_RE = re.compile(r'\{\{[^{}]+\}\}')

    def defragmenter_paragraphe(p_elem) -> bool:
        """Fusionne les runs couvrant une balise fragmentée. Retourne True si modifié."""
        modified = False
        for _ in range(50):  # garde-fou anti-boucle infinie
            runs  = [c for c in p_elem if c.tag == R_TAG]
            texts = []
            for r in runs:
                t = r.find(T_TAG)
                texts.append(t.text if (t is not None and t.text) else '')
            concat = ''.join(texts)
            if '{{' not in concat:
                return modified

            # Offsets de début de chaque run dans le texte concaténé
            offsets = []
            pos = 0
            for txt in texts:
                offsets.append(pos)
                pos += len(txt)

            def run_index(char_pos):
                for i in range(len(offsets) - 1, -1, -1):
                    if offsets[i] <= char_pos:
                        return i
                return 0

            # Chercher la première balise à cheval sur plusieurs runs
            fusion = None
            for m in TAG_SPAN_RE.finditer(concat):
                i = run_index(m.start())
                j = run_index(m.end() - 1)
                if j > i:
                    fusion = (i, j)
                    break
            if fusion is None:
                return modified

            i, j = fusion
            # Vérifier que les runs i..j sont des siblings contigus (pas de <a:br> entre eux)
            children = list(p_elem)
            idx_i = children.index(runs[i])
            idx_j = children.index(runs[j])
            segment = children[idx_i:idx_j + 1]
            if any(c.tag != R_TAG for c in segment):
                return modified  # séparés par un saut de ligne → ne pas fusionner

            # Fusion : concaténer les textes dans le run i, supprimer les runs i+1..j
            merged_text = ''.join(texts[i:j + 1])
            t_elem = runs[i].find(T_TAG)
            if t_elem is None:
                t_elem = etree.SubElement(runs[i], T_TAG)
            t_elem.text = merged_text
            t_elem.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
            for r in runs[i + 1:j + 1]:
                p_elem.remove(r)
            modified = True
        return modified

    src = _io.BytesIO(pptx_bytes)
    out = _io.BytesIO()
    with zipfile.ZipFile(src, 'r') as zin, \
         zipfile.ZipFile(out, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if re.match(r'ppt/slides/slide\d+\.xml$', item.filename):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for p_elem in root.iter(P_TAG):
                        if defragmenter_paragraphe(p_elem):
                            changed = True
                        # Normaliser l'ordre OOXML : endParaRPr doit être le
                        # DERNIER enfant du paragraphe (sinon PowerPoint ignore
                        # les runs placés après — texte invisible au téléchargement).
                        end = p_elem.find(f'{{{A_NS}}}endParaRPr')
                        if end is not None and list(p_elem)[-1] is not end:
                            p_elem.remove(end)
                            p_elem.append(end)
                            changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True,
                                              encoding='UTF-8', standalone=True)
                except Exception:
                    pass  # en cas d'erreur, garder le XML original
            zout.writestr(item, data)
    return out.getvalue()


def extraire_balises_template(template_path: str) -> set:
    """Retourne l'ensemble des balises {{...}} présentes dans les slides du template."""
    TAG_RE_LOCAL = re.compile(r'\{\{([^{}]+)\}\}')
    tags = set()
    with zipfile.ZipFile(template_path) as z:
        for name in z.namelist():
            if re.match(r'ppt/slides/slide\d+\.xml$', name):
                content = z.read(name).decode('utf-8', errors='ignore')
                flat = re.sub(r'<[^>]+>', '', content)
                for t in TAG_RE_LOCAL.findall(flat):
                    tags.add("{{" + t + "}}")
    return tags


def extraire_contextes_balises(template_path: str, fenetre: int = 5) -> dict:
    """
    Pour chaque balise du template, extrait quelques mots avant/après sa
    première occurrence — donné à l'IA pour produire des fragments cohérents.
    """
    TAG_RE_LOCAL = re.compile(r'\{\{[^{}]+\}\}')
    contextes = {}
    with zipfile.ZipFile(template_path) as z:
        for name in sorted(z.namelist()):
            if not re.match(r'ppt/slides/slide\d+\.xml$', name):
                continue
            content = z.read(name).decode('utf-8', errors='ignore')
            # Reconstituer le texte par paragraphe
            for p_xml in re.findall(r'<a:p(?:\s[^>]*)?>.*?</a:p>', content, re.DOTALL):
                texts = re.findall(r'<a:t[^>]*>(.*?)</a:t>', p_xml, re.DOTALL)
                flat = ''.join(texts)
                for m in TAG_RE_LOCAL.finditer(flat):
                    tag = m.group(0)
                    if tag in contextes:
                        continue
                    avant = ' '.join(flat[:m.start()].split()[-fenetre:])
                    apres = ' '.join(flat[m.end():].split()[:fenetre])
                    if avant or apres:
                        contextes[tag] = f"…{avant} ⟦BALISE⟧ {apres}…"
    return contextes


# ============================================================
# === AJUSTEMENT HAUTEUR DES BOÎTES AUTO-FIT =================
# ============================================================

def ajuster_hauteur_autofit(pptx_bytes: bytes) -> bytes:
    """
    Agrandit les boîtes de texte 'auto-ajustées' (spAutoFit) dont la hauteur
    stockée est trop petite pour le contenu généré.

    Pourquoi : LibreOffice recalcule spAutoFit à l'ouverture (donc l'aperçu est
    correct), mais PowerPoint fait confiance à la hauteur stockée et masque le
    texte qui déborde. On pré-dimensionne donc les boîtes pour garantir un
    rendu identique dans tous les logiciels. Plafonné pour ne pas chevaucher la
    shape située en dessous.
    """
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE
    import io as _io

    prs = Presentation(_io.BytesIO(pptx_bytes))
    SH    = prs.slide_height
    MARGE = int(0.15 * 360000)

    for slide in prs.slides:
        boxes = [s for s in slide.shapes
                 if s.has_text_frame and s.top is not None and s.height is not None
                 and s.width and s.left is not None]
        for shape in boxes:
            tf = shape.text_frame
            try:
                if tf.auto_size != MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
                    continue
            except Exception:
                continue

            needed = 0
            for para in tf.paragraphs:
                sz = 14
                for run in para.runs:
                    if run.font.size:
                        sz = run.font.size.pt
                        break
                char_w = sz * 0.5 * 12700
                cpl    = max(1, int(shape.width / char_w))
                txt    = para.text
                nlines = max(1, -(-len(txt) // cpl)) if txt.strip() else 1
                needed += int(nlines * sz * 1.25 * 12700)
            needed += MARGE

            if needed <= shape.height:
                continue

            cap = SH - shape.top - MARGE
            for other in boxes:
                if other is shape or other.top <= shape.top:
                    continue
                # chevauchement horizontal ?
                if other.left < shape.left + shape.width and shape.left < other.left + other.width:
                    cap = min(cap, other.top - shape.top - MARGE)

            new_h = min(needed, cap)
            if new_h > shape.height:
                shape.height = new_h

    out = _io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ============================================================
# === SUPPRESSION DE SLIDES ==================================
# ============================================================

def supprimer_slides(pptx_bytes: bytes, indices_0based: list) -> bytes:
    """
    Supprime les slides aux indices 0-based donnés.
    Travaille sur les bytes du PPTX, retourne les bytes modifiés.
    """
    from pptx import Presentation
    import io as _io

    if not indices_0based:
        return pptx_bytes

    prs = Presentation(_io.BytesIO(pptx_bytes))
    n   = len(prs.slides)
    REL = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'

    for idx in sorted(set(i for i in indices_0based if 0 <= i < n), reverse=True):
        xml_slides = prs.slides._sldIdLst
        slide_elem = xml_slides[idx]
        rId        = slide_elem.get(REL)
        prs.part.drop_rel(rId)
        xml_slides.remove(slide_elem)

    out = _io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ============================================================
# === RÉGÉNÉRATION ET MODIFICATION ===========================
# ============================================================

def regenerer_depuis_replacements(all_replacements: dict, template_path: str) -> bytes:
    """Régénère le PPTX propre depuis le template avec les replacements donnés."""
    work_dir = tempfile.mkdtemp(prefix="propale_")
    try:
        with zipfile.ZipFile(template_path, 'r') as z:
            z.extractall(work_dir)
        preparer_slides_echeancier(work_dir)
        renommer_tags_dans_slides(work_dir)
        remplacer_tout(work_dir, all_replacements)

        output_path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
        with zipfile.ZipFile(output_path, 'w', compression=zipfile.ZIP_DEFLATED) as z:
            for root, dirs, files in os.walk(work_dir):
                for file in sorted(files):
                    fp = os.path.join(root, file)
                    z.write(fp, os.path.relpath(fp, work_dir))
        with open(output_path, 'rb') as f:
            result = f.read()
        os.unlink(output_path)
        return ajuster_hauteur_autofit(result)
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


def interpreter_modification(
    request: str,
    replacements: dict,
    api_key: str,
    tag_index: dict = None,
) -> dict:
    """
    Interprète une demande de modification.
    Si la requête commence par '[N]' ou 'balise N', le tag est résolu directement
    depuis tag_index sans ambiguïté.
    Retourne {"tag": "{{...}}", "new_value": "...", "error": None}
    ou       {"tag": None,      "new_value": None,  "error": "message"}
    """
    client = OpenAI(api_key=api_key, base_url="https://api.groq.com/openai/v1")

    # ── Détection du préfixe numérique [N] / balise N ────────
    m_num = re.match(
        r'(?:balise\s*(?:n°?\s*)?|#)?\s*[\[\(](\d+)[\]\)]\s*[:\-–]?\s*(.*)',
        request.strip(), re.IGNORECASE | re.DOTALL,
    )
    if not m_num:
        # Essai sans crochets : "balise 3 : ..."
        m_num = re.match(
            r'(?:balise\s+(?:n°?\s*)?)(\d+)\s*[:\-–]?\s*(.*)',
            request.strip(), re.IGNORECASE | re.DOTALL,
        )

    if m_num and tag_index:
        num         = int(m_num.group(1))
        instruction = m_num.group(2).strip()
        if num in tag_index:
            tag     = tag_index[num]
            old_val = replacements.get(tag, "")
            tag_nom = tag.replace("{{", "").replace("}}", "")

            prompt = f"""Champ : {tag_nom}
Valeur actuelle : {old_val[:200]}
Instruction : {instruction}

Donne UNIQUEMENT la nouvelle valeur pour ce champ (sans explication, sans guillemets)."""

            resp = client.chat.completions.create(
                model=GROQ_MODEL,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.0,
                max_tokens=300,
            )
            new_val = resp.choices[0].message.content.strip()
            return {"tag": tag, "new_value": new_val, "error": None}
        else:
            return {"tag": None, "new_value": None, "error": f"Balise [{num}] introuvable dans l'index."}

    # ── Fallback : recherche par IA sur tous les champs ───────
    champs_lines = []
    for tag, val in sorted(replacements.items()):
        if val:
            nom   = tag.replace("{{", "").replace("}}", "")
            label = f"{nom}: [À COMPLÉTER]" if "[À COMPLÉTER]" in str(val) else f"{nom}: {str(val)[:120]}"
            champs_lines.append(label)

    champs = "\n".join(champs_lines[:80])

    prompt = f"""Tu aides à modifier une proposition commerciale.

Champs modifiables (nom: valeur actuelle) :
{champs}

Demande : "{request}"

Identifie le champ à modifier et la nouvelle valeur.
Réponds EXACTEMENT dans ce format (rien d'autre) :
TAG: nom_du_champ
VALEUR: la nouvelle valeur

Si aucun champ ne correspond :
ERREUR: explication courte en français"""

    resp = client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.0,
        max_tokens=150,
    )
    raw = resp.choices[0].message.content.strip()

    if raw.startswith("ERREUR:"):
        return {"tag": None, "new_value": None, "error": raw[7:].strip()}

    tag_m = re.search(r'TAG:\s*(\S+)', raw)
    val_m = re.search(r'VALEUR:\s*(.+)', raw, re.DOTALL)

    if tag_m and val_m:
        tag_name = tag_m.group(1).strip().strip("{}")
        new_val  = val_m.group(1).strip()
        tag      = "{{" + tag_name + "}}"
        if tag in replacements:
            return {"tag": tag, "new_value": new_val, "error": None}
        for existing in replacements:
            if tag_name.lower() in existing.lower():
                return {"tag": existing, "new_value": new_val, "error": None}
        return {"tag": None, "new_value": None, "error": f"Champ `{tag_name}` introuvable."}

    return {"tag": None, "new_value": None, "error": "Je n'ai pas compris la demande. Précisez ou utilisez '[N] : nouvelle valeur'."}


def preparer_slides_echeancier(work_dir: str):
    pass  # balises déjà présentes dans la template


def renommer_tags_dans_slides(work_dir: str):
    for opt_num, slide_names in SLIDES_PAR_OPTION.items():
        suffix = f"_opt{opt_num}"
        for slide_name in slide_names:
            slide_path = os.path.join(work_dir, "ppt", "slides", slide_name)
            if not os.path.exists(slide_path):
                continue
            with open(slide_path, "r", encoding="utf-8") as f:
                content = f.read()
            for tag_base in TAGS_FINANCIERS_PAR_OPTION:
                old_tag = "{{" + tag_base + "}}"
                new_tag = "{{" + tag_base + suffix + "}}"
                content = content.replace(old_tag, new_tag)
            with open(slide_path, "w", encoding="utf-8") as f:
                f.write(content)


# ============================================================
# === FONCTION PRINCIPALE ====================================
# ============================================================

def generer_propale(
    resume: str,
    det_paths: list,
    api_key: str,
    template_path: str,
    progress_callback=None,
) -> tuple:
    """
    Génère la propale PPTX.

    Args:
        resume: Texte de résumé du projet
        det_paths: Liste de 3 éléments (chemin fichier ou None)
        api_key: Clé API Groq
        template_path: Chemin absolu vers le template .pptx
        progress_callback: callable(step: int, total: int, message: str)

    Returns:
        (pptx_bytes, all_replacements, missed_tags, nom_client)
    """
    def progress(step, total, msg):
        if progress_callback:
            progress_callback(step, total, msg)

    progress(0, 5, "Lecture des DETs...")

    # Lecture des DETs
    dets_data = []
    for i, chemin in enumerate(det_paths, start=1):
        if chemin and os.path.exists(chemin):
            data = lire_det(chemin)
            dets_data.append((i, data))
        else:
            dets_data.append((i, None))

    dets_valides = [(i, d) for i, d in dets_data if d is not None]

    # Remplacements directs depuis DETs
    progress(1, 5, "Construction des données financières...")
    all_replacements = {}

    det1 = dict(dets_valides).get(1)
    det2 = dict(dets_valides).get(2)
    det3 = dict(dets_valides).get(3)

    if det1:
        all_replacements.update(construire_replacements_directs(det1, suffix=""))
    if det2:
        all_replacements.update(construire_replacements_directs(det2, suffix="_opt2"))
    if det3:
        all_replacements.update(construire_replacements_directs(det3, suffix="_opt3"))

    nb_dets = len(dets_valides)
    all_replacements["{{nb_option}}"] = str(nb_dets) if nb_dets > 0 else "1"

    if dets_valides:
        valid_dets = [d for _, d in dets_valides]
        inter_vals = [str(int(d["nb_intervenants_total"])) for d in valid_dets if d["nb_intervenants_total"]]
        all_replacements["{{nb_intervenants}}"]       = formater_ou_multi(inter_vals)
        all_replacements["{{nb_intervenants_total}}"] = formater_ou_multi(inter_vals)
        all_replacements["{{nb_intervenants_tota}}"]  = formater_ou_multi(inter_vals)
        all_replacements["{{intervenant}}"]           = formater_ou_multi(inter_vals, " intervenants")

        aines = [nombre_en_aine(int(d["nb_qs_phase1"])) for d in valid_dets if d["nb_qs_phase1"]]
        if aines:
            prefix = "une " + formater_ou_multi(aines)
            all_replacements["{{info_questionnaire}}"] = prefix + " de données à majorité quantitative"
            all_replacements["{{of_questionnaire}}"]   = prefix + " de questions à majorité fermées"

        qs_vals = [str(int(d["nb_qs_phase1"])) for d in valid_dets if d["nb_qs_phase1"]]
        all_replacements["{{nb_questionnaire}}"] = formater_ou_multi(qs_vals)

        if len(valid_dets) == 1:
            d = valid_dets[0]
            all_replacements["{{nb_questionnaire_min}}"] = str(d["nb_qs_phase1"]) if d["nb_qs_phase1"] else "[À COMPLÉTER]"
            all_replacements["{{nb_questionnaire_max}}"] = str(d["nb_qs_phase2"]) if d["nb_qs_phase2"] else "[À COMPLÉTER]"
        else:
            pairs = []
            for d in valid_dets:
                vmin = str(int(d["nb_qs_phase1"])) if d["nb_qs_phase1"] else "?"
                vmax = str(int(d["nb_qs_phase2"])) if d["nb_qs_phase2"] else "?"
                pairs.append(f"{vmin},{vmax}")
            all_replacements["{{nb_questionnaire_min}}"] = pairs[0] if pairs else "[À COMPLÉTER]"
            all_replacements["{{nb_questionnaire_max}}"] = " ou ".join(pairs[1:]) if len(pairs) > 1 else "[À COMPLÉTER]"

        listing_vals = [str(int(d["nb_qs_phase1"]) * 2) for d in valid_dets if d["nb_qs_phase1"]]
        all_replacements["{{nb_listing}}"] = formater_ou_multi(listing_vals)

        questio_vals = []
        for d in valid_dets:
            qs   = int(d["nb_qs_phase1"])          if d["nb_qs_phase1"]          else 0
            ints = int(d["nb_intervenants_total"])  if d["nb_intervenants_total"] else 0
            if qs > 0 and ints > 0:
                questio_vals.append(str(_math.ceil(qs / ints)))
        if questio_vals:
            all_replacements["{{questioprinter}}"] = formater_ou_multi(questio_vals)

    # Appel IA (avec contexte d'insertion extrait du template pour les fragments)
    progress(2, 5, "Génération du contenu rédactionnel (IA)...")
    brief = construire_brief_depuis_resume(resume, [d for _, d in dets_valides])
    contextes = extraire_contextes_balises(template_path)
    replacements_ia = appeler_ia(brief, api_key, contextes=contextes)

    if "{{titre_etude}}" in replacements_ia:
        replacements_ia["{{titre_etude}}"] = replacements_ia["{{titre_etude}}"].upper()
    for upper_tag in UPPER_TAGS:
        if upper_tag in replacements_ia:
            replacements_ia[upper_tag] = replacements_ia[upper_tag].upper()

    all_replacements = {**replacements_ia, **all_replacements}

    # Retirer les amorces de phrase répétées par l'IA ('le client souhaite le client souhaite…')
    nettoyer_fragments(all_replacements)

    # Compléter les montants en toutes lettres ({{...ltr}}) depuis les montants numériques
    completer_montants_en_lettres(all_replacements)

    # ── Filet de sécurité : AUCUNE balise brute ne doit rester dans le rendu ──
    # Toute balise du template non couverte reçoit [À COMPLÉTER] (colorée + numérotée,
    # donc modifiable via le chat comme les autres).
    for tag in extraire_balises_template(template_path):
        all_replacements.setdefault(tag, "[À COMPLÉTER]")

    # Génération du PPTX
    progress(3, 5, "Injection des données dans le PowerPoint...")
    work_dir = tempfile.mkdtemp(prefix="propale_")
    try:
        with zipfile.ZipFile(template_path, 'r') as z:
            z.extractall(work_dir)

        preparer_slides_echeancier(work_dir)
        renommer_tags_dans_slides(work_dir)
        replaced = remplacer_tout(work_dir, all_replacements)

        progress(4, 5, "Finalisation du fichier...")
        output_buffer = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        output_path = output_buffer.name
        output_buffer.close()

        with zipfile.ZipFile(output_path, 'w', compression=zipfile.ZIP_DEFLATED) as z:
            for root, dirs, files in os.walk(work_dir):
                for file in sorted(files):
                    fp = os.path.join(root, file)
                    z.write(fp, os.path.relpath(fp, work_dir))

        with open(output_path, 'rb') as f:
            pptx_bytes = f.read()
        os.unlink(output_path)

    finally:
        shutil.rmtree(work_dir, ignore_errors=True)

    # Pré-dimensionner les boîtes auto-ajustées pour un rendu identique dans PowerPoint
    pptx_bytes = ajuster_hauteur_autofit(pptx_bytes)

    missed_tags = set(all_replacements.keys()) - replaced
    nom_client = ""
    if det1 and det1.get("nom_entreprise"):
        nom_client = det1["nom_entreprise"]
    elif dets_valides:
        nom_client = dets_valides[0][1].get("nom_entreprise", "")
    if not nom_client:
        # Utiliser le nom d'entreprise généré par l'IA s'il est exploitable
        ia_nom = all_replacements.get("{{nom_entreprise}}", "")
        if ia_nom and "[À COMPLÉTER]" not in ia_nom:
            nom_client = ia_nom[:60]
    if not nom_client:
        mots = re.findall(r'\w+', resume)[:3]
        nom_client = " ".join(mots) if mots else "client"

    progress(5, 5, "Terminé !")
    return pptx_bytes, all_replacements, missed_tags, nom_client


# ============================================================
# === APERÇU SLIDES ==========================================
# ============================================================

def extraire_contenu_slides(pptx_bytes: bytes) -> list:
    """Extrait le texte de chaque slide pour l'aperçu."""
    try:
        from pptx import Presentation
        import io as _io
        prs = Presentation(_io.BytesIO(pptx_bytes))
    except Exception:
        return []

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and "{{" not in text:
                        texts.append(text)
        slides.append({"num": i, "texts": texts})
    return slides
